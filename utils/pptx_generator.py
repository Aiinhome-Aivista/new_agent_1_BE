import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def safe_text(val):
    if isinstance(val, list):
        return ", ".join(str(v) for v in val)
    return str(val) if val is not None else ""


import re


def replace_tech_labels_with_icons(mermaid_code):
    tech_icons = {
        "react": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/react/react-original.svg",
        "angular": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/angularjs/angularjs-original.svg",
        "vue": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/vuejs/vuejs-original.svg",
        "flask": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/flask/flask-original.svg",
        "django": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/django/django-plain.svg",
        "express": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/express/express-original.svg",
        "nestjs": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/nestjs/nestjs-plain.svg",
        "spring": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/spring/spring-original.svg",
        "go /": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/go/go-original.svg",
        "golang": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/go/go-original.svg",
        "mysql": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/mysql/mysql-original.svg",
        "postgresql": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/postgresql/postgresql-original.svg",
        "postgres": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/postgresql/postgresql-original.svg",
        "mongodb": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/mongodb/mongodb-original.svg",
        "arangodb": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/arangodb.svg",
        "sqlite": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/sqlite/sqlite-original.svg",
        "redis": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/redis/redis-original.svg",
        "azure": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/azure/azure-original.svg",
        "aws": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/amazonaws.svg",
        "google cloud": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/googlecloud.svg",
        "salesforce": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/salesforce.svg",
        "sap": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/sap.svg",
        "tableau": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/tableau.svg",
        "power bi": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/powerbi.svg",
        "powerbi": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/powerbi.svg",
        "github": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/github.svg",
        "devops": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/azuredevops.svg",
        "terraform": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/terraform.svg",
        "stripe": "https://cdn.jsdelivr.net/npm/simple-icons@v9/icons/stripe.svg",
        "docker": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/docker/docker-original.svg",
        "kubernetes": "https://cdn.jsdelivr.net/gh/devicons/devicon/icons/kubernetes/kubernetes-plain.svg"
    }

    pattern = r'(\b[a-zA-Z0-9_-]+\b\s*(?:\[\("?|\["?|\("?|\{"?))([^"\]\)\}]+)("?(?:\)\)|\]\)|\}\}|\]|\)))'
    
    def replacer(match):
        prefix = match.group(1)
        text = match.group(2)
        suffix = match.group(3)
        
        if "<img" in text:
            return match.group(0)
            
        text_lower = text.lower()
        for tech, url in tech_icons.items():
            if tech in text_lower:
                clean_text = text.replace("'", "\\'")
                img_tag = f"<img src='{url}' width='28' height='28'/> "
                has_quotes = prefix.endswith('"') or suffix.startswith('"')
                if not has_quotes:
                    new_prefix = prefix + '"'
                    new_suffix = '"' + suffix
                    return f"{new_prefix}{img_tag}{clean_text}{new_suffix}"
                else:
                    return f"{prefix}{img_tag}{clean_text}{suffix}"
                    
        return match.group(0)
        
    return re.sub(pattern, replacer, mermaid_code)


def crop_whitespace(image_path):
    from PIL import Image, ImageChops
    try:
        with Image.open(image_path) as img:
            # Convert to RGBA
            img_rgba = img.convert('RGBA')
            
            # Extract alpha channel
            alpha = img_rgba.split()[3]
            bbox_alpha = alpha.getbbox()
            
            # Check difference against pure white background
            bg = Image.new('RGBA', img_rgba.size, (255, 255, 255, 255))
            diff = ImageChops.difference(img_rgba, bg)
            bbox_diff = diff.getbbox()
            
            # Select the most appropriate bounding box
            bbox = None
            if bbox_alpha:
                # If there are transparent areas, getextrema returns min < 255
                extrema = alpha.getextrema()
                if extrema[0] < 255:
                    bbox = bbox_alpha
            
            if not bbox:
                bbox = bbox_diff
                
            if bbox:
                # Add 10px margin to avoid tight clipping of borders
                left, top, right, bottom = bbox
                left = max(0, left - 10)
                top = max(0, top - 10)
                right = min(img_rgba.width, right + 10)
                bottom = min(img_rgba.height, bottom + 10)
                
                cropped = img_rgba.crop((left, top, right, bottom))
                cropped.save(image_path)
    except Exception as e:
        print(f"Error cropping image whitespace: {e}")


def render_mermaid_to_image(mermaid_code):
    import base64
    import requests
    import tempfile
    
    if not mermaid_code or not isinstance(mermaid_code, str):
        return None
        
    mermaid_code = mermaid_code.strip()
    
    if mermaid_code.startswith("```mermaid"):
        mermaid_code = mermaid_code[10:]
    elif mermaid_code.startswith("```"):
        mermaid_code = mermaid_code[3:]
    if mermaid_code.endswith("```"):
        mermaid_code = mermaid_code[:-3]
    mermaid_code = mermaid_code.strip()

    # Automatically transform text labels to embed official technology icons
    try:
        mermaid_code = replace_tech_labels_with_icons(mermaid_code)
    except Exception as e:
        print(f"Failed to replace tech labels with brand icons: {e}")
    
    if mermaid_code.startswith("```mermaid"):
        mermaid_code = mermaid_code[10:]
    elif mermaid_code.startswith("```"):
        mermaid_code = mermaid_code[3:]
    if mermaid_code.endswith("```"):
        mermaid_code = mermaid_code[:-3]
    mermaid_code = mermaid_code.strip()
    
    # Strip any custom inline style/class declarations to enforce corporate base theme
    lines = mermaid_code.split('\n')
    cleaned_lines = []
    for line in lines:
        stripped = line.strip()
        if (stripped.startswith("style ") or 
            stripped.startswith("classDef ") or 
            stripped.startswith("class ") or 
            stripped.startswith("linkStyle ")):
            continue
        cleaned_lines.append(line)
    mermaid_code = "\n".join(cleaned_lines)
    
    if "%%{init:" not in mermaid_code:
        style_init = (
            "%%{init: {\n"
            "    'theme': 'base',\n"
            "    'themeVariables': {\n"
            "        'primaryColor': '#ffffff',\n"
            "        'primaryTextColor': '#2d2d2d',\n"
            "        'primaryBorderColor': '#d04a02',\n"
            "        'lineColor': '#4a4a4a',\n"
            "        'secondaryColor': '#f4f6f8',\n"
            "        'tertiaryColor': '#ffffff',\n"
            "        'mainBkg': '#ffffff',\n"
            "        'nodeBorder': '#d04a02',\n"
            "        'clusterBkg': '#f8f9fa',\n"
            "        'clusterBorder': '#cccccc',\n"
            "        'fontSize': '18px',\n"
            "        'subgraphFontSize': '20px',\n"
            "        'labelFontSize': '15px'\n"
            "    }\n"
            "}}%%\n"
        )
        mermaid_code = style_init + mermaid_code
        
    try:
        diagram_bytes = mermaid_code.encode('utf-8')
        urlsafe_base64_bytes = base64.urlsafe_b64encode(diagram_bytes)
        urlsafe_base64_string = urlsafe_base64_bytes.decode('utf-8').replace('=', '')
        
        url = f"https://mermaid.ink/img/{urlsafe_base64_string}"
        response = requests.get(url, timeout=15)
        if response.status_code == 200:
            temp_dir = tempfile.gettempdir()
            temp_path = os.path.join(temp_dir, f"mermaid_{os.urandom(8).hex()}.png")
            with open(temp_path, "wb") as f:
                f.write(response.content)
            
            # Auto-crop white/transparent margins around diagram elements
            crop_whitespace(temp_path)
            
            return temp_path
        else:
            print(f"Mermaid rendering API returned status: {response.status_code}")
    except Exception as e:
        print(f"Error rendering mermaid diagram: {e}")
    return None


def add_picture_proportionally(slide, image_path, left_limit, top_limit, max_w, max_h):
    from PIL import Image
    try:
        with Image.open(image_path) as img:
            img_w, img_h = img.size
        img_ratio = img_w / img_h
        box_ratio = max_w / max_h
        
        if img_ratio > box_ratio:
            fit_w = max_w
            fit_h = max_w / img_ratio
        else:
            fit_h = max_h
            fit_w = max_h * img_ratio
            
        center_x = left_limit + (max_w - fit_w) / 2
        center_y = top_limit + (max_h - fit_h) / 2
        
        slide.shapes.add_picture(image_path, center_x, center_y, fit_w, fit_h)
    except Exception as e:
        print(f"Failed to add picture proportionally: {e}")
        slide.shapes.add_picture(image_path, left_limit, top_limit, max_w, max_h)



# Brand Colors
ORANGE = RGBColor(208, 74, 2)       # #D04A02
CHARCOAL = RGBColor(45, 45, 45)     # #2D2D2D
GOLD = RGBColor(235, 163, 0)        # #EBA300
RED = RGBColor(163, 31, 52)         # #A31F34
OFF_WHITE = RGBColor(248, 249, 250) # #F8F9FA
WHITE = RGBColor(255, 255, 255)
GREY = RGBColor(180, 180, 180)
LIGHT_GREY = RGBColor(240, 240, 240)

def set_font(run, name="Arial", size=14, bold=False, italic=False, color=CHARCOAL):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

def create_slide_header(slide, title_text, subtitle_text=None):
    # Add a top header bar block
    header_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = CHARCOAL
    header_box.line.color.rgb = ORANGE
    header_box.line.width = Pt(1.5)

    # Add title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    
    # Scale font size based on length
    title_size = 24
    if len(title_text) > 70:
        title_size = 14
    elif len(title_text) > 50:
        title_size = 18
        
    set_font(run, size=title_size, bold=True, color=WHITE)

    # Add subtitle
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        run2 = p2.runs[0]
        set_font(run2, size=9, color=GOLD)

def add_footer(slide):
    # Footnote bar
    footer_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(10), Inches(0.4)
    )
    footer_box.fill.solid()
    footer_box.fill.fore_color.rgb = LIGHT_GREY
    footer_box.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(9), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Solution Advisory  |  Autonomous Bid Lifecycle Platform  |  AI Draft - For Internal Review Only"
    p.alignment = PP_ALIGN.LEFT
    set_font(p.runs[0], size=8, bold=False, color=CHARCOAL)


def add_reference_architecture_slide(slide, prs, data):
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    create_slide_header(slide, "Reference Architecture", "Enterprise logical data patterns")
    add_footer(slide)

    cloud_name = "Azure"
    if data and ("aws" in str(data).lower() or "amazon" in str(data).lower()):
        cloud_name = "AWS"
    elif data and ("gcp" in str(data).lower() or "google" in str(data).lower()):
        cloud_name = "Google Cloud"

    storage_name = "Azure Blob Storage"
    if cloud_name == "AWS":
        storage_name = "AWS S3"
    elif cloud_name == "Google Cloud":
        storage_name = "Google Cloud Storage"

    mermaid_code = (
        "graph LR\n"
        "    subgraph External [External Sources]\n"
        "        E1[RFI/RFP Documents]\n"
        "        E2[Project Timeline]\n"
        "        E3[Competency Documents]\n"
        "        E4[Company Asset Lists]\n"
        "        E5[Questionnaires]\n"
        "        E6[Financial Documents]\n"
        "        PA[Parsing Agent]\n"
        "    end\n"
        "    \n"
        "    subgraph Orchestration [Orchestration Layer]\n"
        "        HO[Hierarchical Orchestrator]\n"
        "        RMA[Requirement Mapping Agent]\n"
        "        SDA[Solution Design Agent]\n"
        "        PLA[Planning Agent]\n"
        "        RA[Rendering Agent]\n"
        "    end\n"
        "    \n"
        "    subgraph Knowledge [Knowledge Layer]\n"
        "        KB[PostgreSQL: Competencies]\n"
        "        VS[ChromaDB: Vector Store]\n"
        "    end\n"
        "    \n"
        "    subgraph Guardrails [Guardrails]\n"
        "        VAL[Guardrails SDK: Validation]\n"
        "        CC[Compliance Check]\n"
        "    end\n"
        "    \n"
        "    subgraph Output [Output Layer]\n"
        "        PRE[PowerPoint Rendering Engine]\n"
        f"        ABS[{storage_name}: Draft Proposals]\n"
        "        HRI[Human Review Interface]\n"
        "        AP[Approved Proposal]\n"
        "    end\n"
        "    \n"
        "    E1 --> PA\n"
        "    E2 --> PA\n"
        "    E3 --> PA\n"
        "    E4 --> PA\n"
        "    E5 --> PA\n"
        "    E6 --> PA\n"
        "    PA --> HO\n"
        "    \n"
        "    HO --> RMA\n"
        "    HO --> SDA\n"
        "    HO --> PLA\n"
        "    HO --> RA\n"
        "    \n"
        "    RMA --> KB\n"
        "    SDA --> KB\n"
        "    SDA --> VS\n"
        "    PLA --> KB\n"
        "    PLA --> VS\n"
        "    RA --> VS\n"
        "    \n"
        "    KB --> VAL\n"
        "    VS --> VAL\n"
        "    VAL --> CC\n"
        "    \n"
        "    HO --> CC\n"
        "    HO --> PRE\n"
        "    PRE --> ABS\n"
        "    ABS --> HRI\n"
        "    HRI --> AP\n"
        "    CC --> HRI\n"
    )

    temp_img_path = render_mermaid_to_image(mermaid_code)
    if temp_img_path and os.path.exists(temp_img_path):
        try:
            add_picture_proportionally(slide, temp_img_path, Inches(0.2), Inches(1.0), Inches(9.6), Inches(6.0))
        except Exception as e:
            print(f"Failed to add mermaid image to Reference Architecture slide: {e}")
        finally:
            try:
                os.remove(temp_img_path)
            except:
                pass


def add_azure_landscape_architecture_slide(slide, prs, data):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    cloud_name = "Azure"
    if data and ("aws" in str(data).lower() or "amazon" in str(data).lower()):
        cloud_name = "AWS"
    elif data and ("gcp" in str(data).lower() or "google" in str(data).lower()):
        cloud_name = "Google Cloud"

    create_slide_header(slide, f"Landscape Architecture ({cloud_name} Cloud Platform)", f"{cloud_name} Native Services & Integration Topology")
    add_footer(slide)
        
    db_name = data.get("db_tech", "PostgreSQL") if data else "PostgreSQL"
    backend_name = data.get("backend_tech", "FastAPI") if data else "FastAPI"
    
    # Check cloud tools
    ingestion_tool = "Azure Data Factory"
    monitor_tool = "Azure Monitor"
    storage_tool = "Azure Blob Storage"
    search_tool = "Azure Cognitive Search"
    k8s_cluster = "AKS Cluster"
    
    if cloud_name == "AWS":
        ingestion_tool = "AWS Glue"
        monitor_tool = "AWS CloudWatch"
        storage_tool = "AWS S3"
        search_tool = "Amazon Kendra"
        k8s_cluster = "EKS Cluster"
    elif cloud_name == "Google Cloud":
        ingestion_tool = "Google Cloud Dataflow"
        monitor_tool = "Google Cloud Monitoring"
        storage_tool = "Google Cloud Storage"
        search_tool = "Google Cloud Vertex AI Search"
        k8s_cluster = "GKE Cluster"

    mermaid_code = (
        "graph LR\n"
        "    subgraph OnPrem [On-Premises]\n"
        f"        Client[Client Artefacts] --> ADF[{ingestion_tool}: Ingestion]\n"
        "    end\n"
        "    \n"
        "    subgraph Monitoring [Monitoring]\n"
        f"        Monitor[{monitor_tool}: Logs & Metrics] --> Alerts[Alerts & Governance]\n"
        "    end\n"
        "    \n"
        "    subgraph DataFlow [Data Flow]\n"
        "        SJSON[Structured JSON Intermediate] --> DPS[Deterministic Proposal Structure] --> PPDraft[PowerPoint Draft] --> HR[Human Review] --> AP[Approved Proposal]\n"
        "    end\n"
        "    \n"
        f"    subgraph Cloud [{cloud_name} Cloud Platform]\n"
        f"        AKS[{k8s_cluster}: Multi-Agent System]\n"
        f"        FastAPI[{backend_name}: Agent Services] --> Gateway[API Gateway: {backend_name}] --> ExtAPI[External APIs: PowerPoint Rendering]\n"
        f"        DB[{db_name}: Knowledge Base] --> Embed[pgvector: Embeddings] --> Dashboard[Human Review Dashboard]\n"
        "        VS[ChromaDB: Vector Store] --> Embed\n"
        f"        Blob[{storage_tool}: Artefacts]\n"
        f"        Search[{search_tool}: Indexing] --> RAG[RAG: Requirement Mapping] --> Design[Solution Design]\n"
        "    end\n"
        "    \n"
        "    ADF --> AKS\n"
        "    ADF --> SJSON\n"
        "    AKS --> Monitor\n"
        "    AKS --> DPS\n"
        "    AKS --> FastAPI\n"
        "    AKS --> DB\n"
        "    AKS --> VS\n"
        "    AKS --> Blob\n"
        "    AKS --> Search\n"
    )

    temp_img_path = render_mermaid_to_image(mermaid_code)
    if temp_img_path and os.path.exists(temp_img_path):
        try:
            add_picture_proportionally(slide, temp_img_path, Inches(0.2), Inches(1.0), Inches(9.6), Inches(6.0))
        except Exception as e:
            print(f"Failed to add mermaid image to Landscape slide: {e}")
        finally:
            try:
                os.remove(temp_img_path)
            except:
                pass


def generate_pptx(data, output_path, template_path=None):
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    # Standard 4:3 is default, let's change to 16:9 widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # SLIDE 1: Title Cover (Dark/Orange Premium Design)
    # ----------------------------------------------------
    # Try to find a blank layout, usually index 6 in default templates, but we fall back to the first available if not
    try:
        blank_slide_layout = prs.slide_layouts[6]
    except IndexError:
        blank_slide_layout = prs.slide_layouts[0]
        
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background color
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = CHARCOAL
    bg.line.fill.background()

    # Brand Accent Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    # Title box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_main = tf.paragraphs[0]
    p_main.text = safe_text(data.get("proposal_title", "Autonomous Solution Design"))
    set_font(p_main.runs[0], size=36, bold=True, color=WHITE)
    
    p_meta = tf.add_paragraph()
    p_meta.text = "\nDraft Date: July 2026"
    set_font(p_meta.runs[0], size=11, color=ORANGE)

    # ----------------------------------------------------
    # SLIDE 1A: Business Summary
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Business Summary", "Executive overview of the proposed solution")
    add_footer(slide)

    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
    tf_sum = summary_box.text_frame
    tf_sum.word_wrap = True
    
    business_summary = data.get("business_summary", "No business summary provided.")
    for paragraph in business_summary.split('\n'):
        if paragraph.strip():
            p_sum = tf_sum.add_paragraph()
            p_sum.text = safe_text(paragraph.strip())
            p_sum.alignment = PP_ALIGN.JUSTIFY
            set_font(p_sum.runs[0], size=14, color=CHARCOAL)
            p_sum.space_after = Pt(14)

    # ----------------------------------------------------
    # SLIDE 2: Client Requirements & Gap Analysis
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Client Requirements & Gap Analysis", "RAG-driven competence matching against RFP requirements")
    add_footer(slide)

    # Requirements list (Left panel)
    
    req_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.2))
    tf_req = req_title.text_frame
    tf_req.word_wrap = True
    p = tf_req.paragraphs[0]
    p.text = "Key Client Requirements:"
    set_font(p.runs[0], size=14, bold=True, color=ORANGE)
    p.space_after = Pt(12)
    
    for req in data.get("requirements", ["No requirements specified"]):
        p_item = tf_req.add_paragraph()
        p_item.text = f"• {safe_text(req)}"
        set_font(p_item.runs[0], size=12, color=CHARCOAL)
        p_item.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 3: Capability Gaps & Mitigations
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Capability Gaps & Mitigations", "Identified gaps against RFP requirements and proposed mitigations")
    add_footer(slide)



    gap_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.2))
    tf_gap = gap_title.text_frame
    tf_gap.word_wrap = True
    p_gap = tf_gap.paragraphs[0]
    p_gap.text = "Capability Gaps & Mitigations:"
    set_font(p_gap.runs[0], size=14, bold=True, color=RED)
    p_gap.space_after = Pt(12)

    for gap in data.get("gaps", ["No gaps identified"]):
        p_item = tf_gap.add_paragraph()
        p_item.text = f"• {safe_text(gap)}"
        set_font(p_item.runs[0], size=12, color=CHARCOAL)
        p_item.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 3A: Capability Gaps & Mitigations (Cont.)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Capability Gaps & Mitigations", "Identified gaps against RFP requirements and proposed mitigations")
    add_footer(slide)

    gap_title_cont = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.2))
    tf_gap_cont = gap_title_cont.text_frame
    tf_gap_cont.word_wrap = True
    p_gap_cont = tf_gap_cont.paragraphs[0]
    p_gap_cont.text = "Capability Gaps & Mitigations:"
    set_font(p_gap_cont.runs[0], size=14, bold=True, color=RED)
    p_gap_cont.space_after = Pt(12)

    extra_gaps = [
        "Identified gap in Client Requirement: 'Security and Data Protection'. Mitigation: Implement enterprise-grade zero-trust architecture, advanced data encryption, RBAC, and continuous threat detection mechanisms.",
        "Identified gap in Client Requirement: 'System Observability and Monitoring'. Mitigation: Deploy comprehensive distributed tracing, centralized logging, and real-time alerting dashboards across the multi-agent environment."
    ]

    for gap in extra_gaps:
        p_item = tf_gap_cont.add_paragraph()
        p_item.text = f"• {safe_text(gap)}"
        set_font(p_item.runs[0], size=12, color=CHARCOAL)
        p_item.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 4: Solution Approach & Architecture
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Solution Approach & Architecture", "High-level implementation strategy and operational frameworks")
    add_footer(slide)

    # 3-Pillar Solution display
    pillars = data.get("solution_pillars", [
        {"title": "Pillar 1", "desc": "Description 1"},
        {"title": "Pillar 2", "desc": "Description 2"},
        {"title": "Pillar 3", "desc": "Description 3"}
    ])

    width_pillar = Inches(2.7)
    gap_pillar = Inches(0.4)
    start_left = Inches(0.5)

    for i, pillar in enumerate(pillars[:3]):
        left = start_left + i * (width_pillar + gap_pillar)
        
        # Pillar Title Box (no background fill)
        t_box = slide.shapes.add_textbox(left, Inches(1.8), width_pillar, Inches(0.6))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.alignment = PP_ALIGN.LEFT
        p_t.text = ""
        
        # Number Run
        r_num = p_t.add_run()
        r_num.text = f"0{i+1}. "
        set_font(r_num, size=14, bold=True, color=ORANGE)
        
        # Title Run
        r_title = p_t.add_run()
        r_title.text = safe_text(pillar.get('title'))
        set_font(r_title, size=14, bold=True, color=CHARCOAL)
        
        # Horizontal Separator Line (Rectangle)
        hline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.5), width_pillar, Inches(0.02))
        hline.fill.solid()
        hline.fill.fore_color.rgb = ORANGE
        hline.line.fill.background()
        
        # Pillar Description text box
        desc_box = slide.shapes.add_textbox(left, Inches(2.6), width_pillar, Inches(3.5))
        tf_desc = desc_box.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = safe_text(pillar.get("desc", ""))
        p_desc.alignment = PP_ALIGN.JUSTIFY
        set_font(p_desc.runs[0], size=10, color=CHARCOAL)

    # ----------------------------------------------------
    # SLIDE 3B: High Level Design: Data Flow (Custom Architecture Diagram)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Black Header
    hdr_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = RGBColor(15, 15, 15)
    hdr_bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = "High Level Design: Data Flow"
    set_font(p_title.runs[0], size=28, bold=True, color=WHITE)

    # Dynamic Data Flow Layers
    data_flow_items = data.get("data_flow", [])
    if data_flow_items:
        num_items = len(data_flow_items)
        
        # Calculate available height
        start_top = 1.5
        end_bottom = 7.0
        total_avail = end_bottom - start_top
        
        arrow_h = 0.3
        
        # Calculate dynamic box height, cap it so it doesn't get too thick
        box_h = (total_avail - ((num_items - 1) * arrow_h)) / num_items
        if box_h > 0.9:
            box_h = 0.9
        if box_h < 0.4:
            box_h = 0.4
            arrow_h = 0.1
            
        total_block_height = (box_h * num_items) + (arrow_h * (num_items - 1))
        current_top = start_top + (total_avail - total_block_height) / 2
        
        box_width = 8.0
        box_left = 1.0
        
        for i, item_text in enumerate(data_flow_items):
            # Draw Box
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box_left), Inches(current_top), Inches(box_width), Inches(box_h))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(30, 30, 30)
            box.line.color.rgb = RGBColor(120, 120, 120)
            box.line.width = Pt(1.5)
            
            # Text inside box
            tb = slide.shapes.add_textbox(Inches(box_left + 0.1), Inches(current_top), Inches(box_width - 0.2), Inches(box_h))
            tf = tb.text_frame
            tf.word_wrap = True
            try:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except:
                pass
            p = tf.paragraphs[0]
            p.text = safe_text(item_text)
            p.alignment = PP_ALIGN.CENTER
            set_font(p.runs[0], size=11, bold=True, color=WHITE)
            
            current_top += box_h
            
            # Draw downward arrow if not the last item
            if i < num_items - 1:
                arr = slide.shapes.add_textbox(Inches(4.8), Inches(current_top), Inches(0.4), Inches(arrow_h))
                pa = arr.text_frame.paragraphs[0]
                pa.text = "▼"
                pa.alignment = PP_ALIGN.CENTER
                set_font(pa.runs[0], size=12, bold=True, color=ORANGE)
                current_top += arrow_h

    # ----------------------------------------------------
    # SLIDE 4: Landscape Architecture
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Landscape & Architecture", "Reference systems architecture and integration patterns")
    add_footer(slide)

    arch_layers = data.get("architecture", [
        {"name": "Client Access / Presentation Layer", "components": ["Web Portal", "Mobile Client", "API Gateway"]},
        {"name": "Application Logic & Agents Core", "components": ["Orchestrator Engine", "Estimation Engine", "Document Agent"]},
        {"name": "Data Integration & Knowledge", "components": ["MySQL Database", "Qdrant Vector DB", "Asset Library"]}
    ])

    current_top = 1.5
    for i, layer in enumerate(arch_layers[:4]):
        comps = layer.get("components", [])
        
        # Max 4 components per row to prevent squashing
        max_per_row = 4
        num_rows = (len(comps) + max_per_row - 1) // max_per_row if len(comps) > 0 else 1
        layer_h = 1.0 + (num_rows - 1) * 0.8
        
        # Layer Container
        layer_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(current_top), Inches(9.0), Inches(layer_h))
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = OFF_WHITE
        layer_box.line.color.rgb = GOLD
        layer_box.line.width = Pt(1.5)
        
        # Layer Header Box (left side)
        hdr_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(current_top), Inches(2.5), Inches(layer_h))
        hdr_box.fill.solid()
        hdr_box.fill.fore_color.rgb = CHARCOAL
        hdr_box.line.fill.background()
        
        p_hdr = hdr_box.text_frame.paragraphs[0]
        p_hdr.text = safe_text(layer.get("name", ""))
        p_hdr.alignment = PP_ALIGN.CENTER
        set_font(p_hdr.runs[0], size=11, bold=True, color=WHITE)
        
        # Add downward arrow if not the last layer
        if i < len(arch_layers[:4]) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.6), Inches(current_top + layer_h + 0.05), Inches(0.3), Inches(0.2))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()
        
        # Component boxes inside the layer
        if comps:
            for r in range(num_rows):
                row_comps = comps[r * max_per_row : (r + 1) * max_per_row]
                c_width = Inches(6.0 / len(row_comps)) if len(row_comps) > 0 else Inches(1.5)
                
                for j, comp in enumerate(row_comps):
                    c_left = Inches(3.2) + (j * c_width)
                    c_top = current_top + 0.15 + (r * 0.8)
                    
                    c_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(c_top), c_width - Inches(0.15), Inches(0.7))
                    c_box.fill.solid()
                    c_box.fill.fore_color.rgb = WHITE
                    c_box.line.color.rgb = ORANGE
                    c_box.line.width = Pt(1)
                    
                    tf_comp = c_box.text_frame
                    tf_comp.word_wrap = True
                    try:
                        tf_comp.vertical_anchor = MSO_ANCHOR.MIDDLE
                    except:
                        pass
                    p_comp = tf_comp.paragraphs[0]
                    p_comp.text = safe_text(comp)
                    p_comp.alignment = PP_ALIGN.CENTER
                    set_font(p_comp.runs[0], size=8.5, bold=True, color=CHARCOAL)
                    
                    # Add horizontal arrow to next component to show sequence
                    if j < len(row_comps) - 1:
                        arr_x = c_left + c_width - Inches(0.15)
                        arr_y = c_top + 0.25
                        c_arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arr_x, Inches(arr_y), Inches(0.15), Inches(0.2))
                        c_arr.fill.solid()
                        c_arr.fill.fore_color.rgb = ORANGE
                        c_arr.line.fill.background()

        current_top += layer_h + 0.3

    # ----------------------------------------------------
    # SLIDE 4B: Azure Cost Calculator (Estimations)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Infrastructure Approximation", "Estimated cloud infrastructure components and costs")
    add_footer(slide)

    infra_items = data.get("infrastructure_approximation", [])
    
    if infra_items:
        infra_rows = len(infra_items) + 1
        infra_cols = 3
        infra_table_shape = slide.shapes.add_table(infra_rows, infra_cols, Inches(1.0), Inches(2.0), Inches(8.0), Inches(0.5 * infra_rows))
        i_table = infra_table_shape.table
        
        i_table.columns[0].width = Inches(2.5)
        i_table.columns[1].width = Inches(3.5)
        i_table.columns[2].width = Inches(2.0)
        
        i_headers = ["Azure Component", "Specification", "Unit Cost"]
        for j, hdr in enumerate(i_headers):
            cell = i_table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CHARCOAL
            p = cell.text_frame.paragraphs[0]
            p.text = hdr
            p.alignment = PP_ALIGN.CENTER
            set_font(p.runs[0], size=12, bold=True, color=WHITE)
            
        for i, item in enumerate(infra_items):
            row_idx = i + 1
            cost_raw = item.get("estimated_monthly_cost")
            cost_str = f"$ {cost_raw} onwards per hour" if cost_raw is not None else ""
            vals = [item.get("component"), item.get("spec"), cost_str]
            for j, val in enumerate(vals):
                cell = i_table.cell(row_idx, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else OFF_WHITE
                p = cell.text_frame.paragraphs[0]
                p.text = safe_text(val)
                p.alignment = PP_ALIGN.CENTER if j == 2 else PP_ALIGN.LEFT
                set_font(p.runs[0], size=11, color=CHARCOAL)

    # ----------------------------------------------------
    # SLIDE 5: Project Milestones
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Project Milestones", "Sequential delivery phases and target deliverables")
    add_footer(slide)

    phases = data.get("timeline_phases", [
        {"phase": "Design & Planning", "deliverables": "RFP requirements analysis"},
        {"phase": "Development", "deliverables": "Core engineering & integration"},
        {"phase": "Testing", "deliverables": "QA and Integration Testing"},
        {"phase": "Deployment", "deliverables": "Production release"},
        {"phase": "Training", "deliverables": "User training & handover"}
    ])

    enforced_names = ["Design & Planning", "Development", "Testing", "Deployment", "Training"]

    for i, phase in enumerate(phases[:5]):
        top_val = 1.35 + (i * 1.12)
        top = Inches(top_val)
        
        # Enforce exact phase name
        if i < len(enforced_names):
            phase_name = enforced_names[i]
        else:
            phase_name = safe_text(phase.get("phase", ""))
        
        # Chevron Phase Box
        c_shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0.5), Inches(top_val + 0.1), Inches(3.2), Inches(0.85))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = ORANGE
        c_shape.line.fill.background()
        
        # Chevron Text
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = phase_name
        p_c.alignment = PP_ALIGN.CENTER
        set_font(p_c.runs[0], size=11, bold=True, color=WHITE)

        # Description / Deliverables Box next to the Chevron
        d_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), top, Inches(5.5), Inches(1.05))
        d_box.fill.solid()
        d_box.fill.fore_color.rgb = OFF_WHITE
        d_box.line.color.rgb = CHARCOAL
        d_box.line.width = Pt(1)
        
        tf_d = d_box.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = "Key Deliverables / Activities:"
        set_font(p_d.runs[0], size=11, bold=True, color=CHARCOAL)
        
        p_d_desc = tf_d.add_paragraph()
        p_d_desc.text = safe_text(phase.get("deliverables", ""))
        set_font(p_d_desc.runs[0], size=9, color=CHARCOAL)

    # ----------------------------------------------------
    # SLIDE 5B: Case Study
    # ----------------------------------------------------
    sim_projects = data.get("similar_projects", [])
    
    if sim_projects:
        for idx_proj, proj in enumerate(sim_projects):
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Dynamic header title
            p_name = safe_text(proj.get('project_name', 'Migration'))
            c_ind = safe_text(proj.get('client_industry', 'Client'))
            title_text = f"Case Study {idx_proj + 1}: {p_name} ({c_ind})"
            create_slide_header(slide, title_text, "Past credentials and successful delivery outcomes")
            add_footer(slide)
            
            # Left Column: High Level Summary Box
            left_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.2))
            left_box.fill.solid()
            left_box.fill.fore_color.rgb = WHITE
            left_box.line.color.rgb = CHARCOAL
            left_box.line.width = Pt(1)
            
            # Label on the border for High Level Summary
            label_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(2.2), Inches(0.3))
            tf_lbl = label_box.text_frame
            tf_lbl.word_wrap = True
            tf_lbl.margin_left = Inches(0.05)
            tf_lbl.margin_right = Inches(0.05)
            tf_lbl.margin_top = Inches(0)
            tf_lbl.margin_bottom = Inches(0)
            p_lbl = tf_lbl.paragraphs[0]
            p_lbl.text = " High Level Summary "
            set_font(p_lbl.runs[0], size=11, bold=True, color=CHARCOAL)
            label_box.fill.solid()
            label_box.fill.fore_color.rgb = WHITE
            label_box.line.fill.background()
            
            # Content inside Left Column
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(4.1), Inches(4.9))
            tf_c = content_box.text_frame
            tf_c.word_wrap = True
            tf_c.margin_left = Inches(0)
            tf_c.margin_right = Inches(0)
            tf_c.margin_top = Inches(0)
            tf_c.margin_bottom = Inches(0)
            
            # Business Problem
            p_bp = tf_c.paragraphs[0]
            p_bp.text = "Business Problem"
            set_font(p_bp.runs[0], size=11, bold=True, color=CHARCOAL)
            p_bp.space_after = Pt(2)
            
            bp_list = proj.get("business_problem", [])
            if not bp_list:
                bp_list = ["No business problem description provided."]
            for bp in bp_list:
                p = tf_c.add_paragraph()
                p.text = f"• {safe_text(bp)}"
                set_font(p.runs[0], size=7.5, color=CHARCOAL)
                p.space_after = Pt(2)
                
            # Spacing
            p_space = tf_c.add_paragraph()
            p_space.text = ""
            p_space.space_after = Pt(2)
            
            # Our Approach
            p_ap = tf_c.add_paragraph()
            p_ap.text = "Our Approach"
            set_font(p_ap.runs[0], size=11, bold=True, color=CHARCOAL)
            p_ap.space_after = Pt(2)
            
            ap_list = proj.get("our_approach", [])
            if not ap_list:
                ap_list = ["No approach description provided."]
            for ap in ap_list:
                p = tf_c.add_paragraph()
                p.text = f"• {safe_text(ap)}"
                set_font(p.runs[0], size=7.5, color=CHARCOAL)
                p.space_after = Pt(2)
            
            # Right Column Top: Technical Architecture Box
            right_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(1.5), Inches(4.5), Inches(3.6))
            right_box.fill.solid()
            right_box.fill.fore_color.rgb = WHITE
            right_box.line.color.rgb = CHARCOAL
            right_box.line.width = Pt(1)
            
            # Label on the border for Technical Architecture
            r_label_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.35), Inches(2.2), Inches(0.3))
            tf_rlbl = r_label_box.text_frame
            tf_rlbl.word_wrap = True
            tf_rlbl.margin_left = Inches(0.05)
            tf_rlbl.margin_right = Inches(0.05)
            tf_rlbl.margin_top = Inches(0)
            tf_rlbl.margin_bottom = Inches(0)
            p_rlbl = tf_rlbl.paragraphs[0]
            p_rlbl.text = " Technical Architecture "
            set_font(p_rlbl.runs[0], size=11, bold=True, color=CHARCOAL)
            r_label_box.fill.solid()
            r_label_box.fill.fore_color.rgb = WHITE
            r_label_box.line.fill.background()
            
            # Mermaid Diagram
            mermaid_code = proj.get("tech_architecture_mermaid")
            if mermaid_code:
                try:
                    temp_img_path = render_mermaid_to_image(mermaid_code)
                    if temp_img_path and os.path.exists(temp_img_path):
                        add_picture_proportionally(slide, temp_img_path, Inches(5.1), Inches(1.8), Inches(4.3), Inches(1.8))
                except Exception as img_e:
                    print(f"Failed to add case study mermaid image: {img_e}")
                    
            # 3 Side-by-Side Explanation boxes underneath the diagram
            explanations = proj.get("tech_architecture_explanation", [])
            while len(explanations) < 3:
                explanations.append("N/A")
                
            box_width = Inches(1.35)
            box_height = Inches(1.3)
            box_top = Inches(3.7)
            
            for idx_exp, exp_text in enumerate(explanations[:3]):
                box_left = Inches(5.1 + idx_exp * 1.45)
                s_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_left, box_top, box_width, box_height)
                s_box.fill.solid()
                s_box.fill.fore_color.rgb = OFF_WHITE
                s_box.line.color.rgb = CHARCOAL
                s_box.line.width = Pt(0.5)
                
                tf_s = s_box.text_frame
                tf_s.word_wrap = True
                tf_s.margin_left = Inches(0.05)
                tf_s.margin_right = Inches(0.05)
                tf_s.margin_top = Inches(0.05)
                tf_s.margin_bottom = Inches(0.05)
                
                p_s = tf_s.paragraphs[0]
                p_s.text = safe_text(exp_text)
                set_font(p_s.runs[0], size=7, color=CHARCOAL)
                
            # Right Column Bottom Left: Key Technologies
            tech_header_box = slide.shapes.add_textbox(Inches(5.0), Inches(5.1), Inches(2.2), Inches(0.3))
            tf_th = tech_header_box.text_frame
            tf_th.word_wrap = True
            p_th = tf_th.paragraphs[0]
            p_th.text = "Key Technologies"
            set_font(p_th.runs[0], size=11, bold=True, color=CHARCOAL)
            
            tech_list_box = slide.shapes.add_textbox(Inches(5.0), Inches(5.4), Inches(2.2), Inches(1.6))
            tf_tl = tech_list_box.text_frame
            tf_tl.word_wrap = True
            tf_tl.margin_left = Inches(0)
            tf_tl.margin_right = Inches(0)
            tf_tl.margin_top = Inches(0)
            tf_tl.margin_bottom = Inches(0)
            
            tech_list = proj.get("key_technologies", [])[:3]
            for tech in tech_list:
                p = tf_tl.add_paragraph() if tf_tl.paragraphs[0].text else tf_tl.paragraphs[0]
                p.text = f"• {safe_text(tech)}"
                set_font(p.runs[0], size=7.0, color=CHARCOAL)
                p.space_after = Pt(1)
                
            # Right Column Bottom Right: Benefits / Outcome
            ben_header_box = slide.shapes.add_textbox(Inches(7.3), Inches(5.1), Inches(2.2), Inches(0.3))
            tf_bh = ben_header_box.text_frame
            tf_bh.word_wrap = True
            p_bh = tf_bh.paragraphs[0]
            p_bh.text = "Benefits/Outcome"
            set_font(p_bh.runs[0], size=11, bold=True, color=CHARCOAL)
            
            ben_list_box = slide.shapes.add_textbox(Inches(7.3), Inches(5.4), Inches(2.2), Inches(1.6))
            tf_bl = ben_list_box.text_frame
            tf_bl.word_wrap = True
            tf_bl.margin_left = Inches(0)
            tf_bl.margin_right = Inches(0)
            tf_bl.margin_top = Inches(0)
            tf_bl.margin_bottom = Inches(0)
            
            ben_list = proj.get("benefits_outcome", [])[:3]
            for ben in ben_list:
                p = tf_bl.add_paragraph() if tf_bl.paragraphs[0].text else tf_bl.paragraphs[0]
                p.text = f"• {safe_text(ben)}"
                set_font(p.runs[0], size=7.0, color=CHARCOAL)
                p.space_after = Pt(1)

    # ----------------------------------------------------
    # SLIDE 6: Effort & Person-Hour Conversion
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Effort & Person-Hour Conversion", "Allocated program FTE structure, rate cards, and financial sizing")
    add_footer(slide)

    # Table layout
    resources = data.get("resources", [
        {"role": "Engagement Partner", "fte": "0.25", "rate": "$30,000", "total": "$45,000", "person_days": 10},
        {"role": "Lead Architect", "fte": "1.00", "rate": "$24,000", "total": "$144,000", "person_days": 60},
        {"role": "Senior Frontend Developer", "fte": "2.00", "rate": "$8,000", "total": "$96,000", "person_days": 120},
        {"role": "Senior Backend Developer", "fte": "2.00", "rate": "$8,000", "total": "$96,000", "person_days": 120},
        {"role": "DevOps & Security Specialist", "fte": "1.00", "rate": "$9,000", "total": "$54,000", "person_days": 60}
    ])
    
    rows = len(resources) + 2  # +1 for header, +1 for Total Assumption
    # Cap rows to fit on one slide
    rows = min(rows, 9)
    cols = 4
    
    # Calculate a sensible height for the table depending on rows
    table_height = Inches(0.4 * rows)
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9.0), table_height)
    table = table_shape.table

    # Column Widths
    table.columns[0].width = Inches(3.5) # Role
    table.columns[1].width = Inches(1.8) # Hourly Rate
    table.columns[2].width = Inches(1.7) # Person Days
    table.columns[3].width = Inches(2.0) # Total Sizing

    headers = ["Role / Competency", "Hourly Rate", "Person Days", "Total Financial Sizing"]
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CHARCOAL
        p = cell.text_frame.paragraphs[0]
        p.text = safe_text(header)
        p.alignment = PP_ALIGN.CENTER
        set_font(p.runs[0], size=11, bold=True, color=WHITE)

    for i, res in enumerate(resources[:rows-2]):
        row_idx = i + 1
        cols_val = [res.get("role"), res.get("rate"), str(res.get("person_days", "N/A")), res.get("total")]
        for j, val in enumerate(cols_val):
            cell = table.cell(row_idx, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else OFF_WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = safe_text(val)
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            set_font(p.runs[0], size=10, bold=(j == 0), color=CHARCOAL)
            
    # Add Total Assumption Row
    last_row_idx = rows - 1
    
    total_hourly_rate = 0
    for res in resources[:rows-2]:
        rate_str = res.get("rate", "0")
        try:
            rate_val = float(str(rate_str).replace('$', '').replace(',', '').strip())
            total_hourly_rate += rate_val
        except:
            pass

    for j in range(cols):
        cell = table.cell(last_row_idx, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if j == 0:
            p.text = "Total Assumption"
            p.alignment = PP_ALIGN.LEFT
        elif j == 1:
            p.text = f"${int(total_hourly_rate):,}"
        elif j == 3:
            p.text = str(safe_text(data.get("budget", "N/A")))
        else:
            p.text = " "
        if p.runs:
            set_font(p.runs[0], size=11, bold=True, color=WHITE)

    # Add Disclaimer below table
    disclaimer_top = Inches(1.5 + (0.4 * rows) + 0.3)
    disclaimer_box = slide.shapes.add_textbox(Inches(0.5), disclaimer_top, Inches(9.0), Inches(0.8))
    tf_disc = disclaimer_box.text_frame
    tf_disc.word_wrap = True
    p_disc = tf_disc.paragraphs[0]
    p_disc.text = "Disclaimer: Please note that this high-level estimate is subject to change as it depends on detailed client requirements. All resource and pricing calculations reflect the median baseline for similar enterprise integrations."
    set_font(p_disc.runs[0], size=10, italic=True, color=CHARCOAL)

    # ----------------------------------------------------
    # SLIDE 7: Required Skills & Competency Matching
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Skills Inventory & Competency Mapping", "Required technical capabilities grounded in organizational assets")
    add_footer(slide)

    skills_map = data.get("skills_mapping", [
        {"skill": "React 18, TypeScript, Tailwind", "role": "Frontend Developer", "conf": "[✔]"},
        {"skill": "Flask API, Python Core", "role": "Backend Developer", "conf": "[✔]"},
        {"skill": "MySQL Connector, RAG Store", "role": "Database Architect", "conf": "[✔]"},
        {"skill": "python-pptx Engine", "role": "Orchestrator Agent", "conf": "[✔]"},
        {"skill": "CI/CD & DevOps", "role": "DevOps Engineer", "conf": "[✔]"}
    ])
    
    rows2 = len(skills_map) + 1
    # Cap rows to fit on one slide
    rows2 = min(rows2, 9)
    cols2 = 2
    
    table_height2 = Inches(0.4 * rows2)
    table_shape2 = slide.shapes.add_table(rows2, cols2, Inches(0.5), Inches(1.5), Inches(9.0), table_height2)
    table2 = table_shape2.table

    table2.columns[0].width = Inches(4.5) # Skill Name
    table2.columns[1].width = Inches(4.5) # Target Role

    headers2 = ["Technical Skill", "Target Project Role"]
    for j, header in enumerate(headers2):
        cell = table2.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        p = cell.text_frame.paragraphs[0]
        p.text = safe_text(header)
        p.alignment = PP_ALIGN.CENTER
        set_font(p.runs[0], size=11, bold=True, color=WHITE)

    for i, item in enumerate(skills_map[:rows2-1]):
        row_idx = i + 1
        cols_val = [item.get("skill"), item.get("role")]
        for j, val in enumerate(cols_val):
            cell = table2.cell(row_idx, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else OFF_WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = safe_text(val)
            p.alignment = PP_ALIGN.LEFT
            set_font(p.runs[0], size=10, bold=(j == 0), color=CHARCOAL)

    # ----------------------------------------------------
    # HARDCODED "SAME TO SAME" ARCHITECTURE DIAGRAMS
    # ----------------------------------------------------
    # Slide 8: Reference Architecture
    ref_slide = prs.slides.add_slide(blank_slide_layout)
    add_reference_architecture_slide(ref_slide, prs, data)

    # Slide 9: Azure Landscape Architecture
    azure_slide = prs.slides.add_slide(blank_slide_layout)
    add_azure_landscape_architecture_slide(azure_slide, prs, data)





    # Slide 10: Thank You Slide
    thank_you_slide = prs.slides.add_slide(blank_slide_layout)
    
    # White background for a clean finish
    bg_ty = thank_you_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg_ty.fill.solid()
    bg_ty.fill.fore_color.rgb = WHITE
    bg_ty.line.fill.background()

    # "Thank You" Main Text
    ty_box = thank_you_slide.shapes.add_textbox(Inches(0), Inches(3.0), Inches(10), Inches(1.5))
    p_ty = ty_box.text_frame.paragraphs[0]
    p_ty.text = "Thank You"
    p_ty.alignment = PP_ALIGN.CENTER
    set_font(p_ty.runs[0], size=72, bold=True, color=ORANGE)

    # Save presentation
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")
