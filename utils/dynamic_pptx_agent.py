import os
import json
from langchain_core.tools import tool
from langchain_openai import ChatOpenAI
from langgraph.prebuilt import create_react_agent
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Import formatting helpers from existing generator
from utils.pptx_generator import (
    safe_text, set_font, create_slide_header, add_footer,
    ORANGE, CHARCOAL, GOLD, RED, OFF_WHITE, WHITE, LIGHT_GREY
)

# Global variables for the agent state
_prs = None
_data = None

@tool
def create_title_slide() -> str:
    """Creates the title/cover slide using project context data."""
    global _prs, _data
    if not _prs or not _data: return "Error: Agent state not initialized."
    
    blank_slide_layout = _prs.slide_layouts[6]
    slide = _prs.slides.add_slide(blank_slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = CHARCOAL
    bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_pre = tf.paragraphs[0]
    p_pre.text = "PWC IT SOLUTION PROPOSAL"
    set_font(p_pre.runs[0], size=14, bold=True, color=GOLD)
    
    p_main = tf.add_paragraph()
    p_main.text = safe_text(_data.get("proposal_title", "Autonomous Solution Design"))
    set_font(p_main.runs[0], size=36, bold=True, color=WHITE)
    
    p_sub = tf.add_paragraph()
    p_sub.text = f"Prepared for: {safe_text(_data.get('client_name', 'Enterprise Client'))}"
    set_font(p_sub.runs[0], size=18, color=LIGHT_GREY)
    
    p_meta = tf.add_paragraph()
    p_meta.text = f"\nTimeline: {safe_text(_data.get('project_duration', 'N/A'))}  |  Target Budget: {safe_text(_data.get('budget', 'N/A'))}\nDraft Date: Current"
    set_font(p_meta.runs[0], size=11, color=ORANGE)
    return "Title slide created successfully."

@tool
def create_business_summary_slide() -> str:
    """Creates the business summary slide."""
    global _prs, _data
    blank_slide_layout = _prs.slide_layouts[6]
    slide = _prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Business Summary", "Executive overview of the proposed solution")
    add_footer(slide)

    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
    tf_sum = summary_box.text_frame
    tf_sum.word_wrap = True
    
    business_summary = _data.get("business_summary", "No business summary provided.")
    for paragraph in business_summary.split('\n'):
        if paragraph.strip():
            p_sum = tf_sum.add_paragraph()
            p_sum.text = safe_text(paragraph.strip())
            set_font(p_sum.runs[0], size=14, color=CHARCOAL)
            p_sum.space_after = Pt(14)
    return "Business summary slide created successfully."

@tool
def create_gap_analysis_slide() -> str:
    """Creates the Client Requirements & Gap Analysis slide."""
    global _prs, _data
    blank_slide_layout = _prs.slide_layouts[6]
    slide = _prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Client Requirements & Gap Analysis", "RAG-driven competence matching against RFP requirements")
    add_footer(slide)

    # Requirements list (Left panel)
    req_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5))
    req_box.fill.solid()
    req_box.fill.fore_color.rgb = OFF_WHITE
    req_box.line.color.rgb = CHARCOAL
    req_box.line.width = Pt(1)
    
    req_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(4.1), Inches(5.2))
    tf_req = req_title.text_frame
    tf_req.word_wrap = True
    p = tf_req.paragraphs[0]
    p.text = "Key Client Requirements:"
    set_font(p.runs[0], size=14, bold=True, color=ORANGE)
    
    for req in _data.get("requirements", ["No requirements specified"]):
        p_item = tf_req.add_paragraph()
        p_item.text = f"• {safe_text(req)}"
        set_font(p_item.runs[0], size=11, color=CHARCOAL)

    # Gaps and Matches (Right panel)
    gap_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.5))
    gap_box.fill.solid()
    gap_box.fill.fore_color.rgb = OFF_WHITE
    gap_box.line.color.rgb = CHARCOAL
    gap_box.line.width = Pt(1)

    gap_title = slide.shapes.add_textbox(Inches(5.3), Inches(1.4), Inches(4.1), Inches(5.2))
    tf_gap = gap_title.text_frame
    tf_gap.word_wrap = True
    p_gap = tf_gap.paragraphs[0]
    p_gap.text = "Capability Gaps & Mitigations:"
    set_font(p_gap.runs[0], size=14, bold=True, color=RED)

    for gap in _data.get("gaps", ["No gaps identified"]):
        p_item = tf_gap.add_paragraph()
        p_item.text = f"• {safe_text(gap)}"
        set_font(p_item.runs[0], size=11, color=CHARCOAL)
    return "Gap analysis slide created successfully."

@tool
def create_timeline_slide() -> str:
    """Creates the Project Timeline slide."""
    global _prs, _data
    blank_slide_layout = _prs.slide_layouts[6]
    slide = _prs.slides.add_slide(blank_slide_layout)
    create_slide_header(slide, "Project Timeline", "Sequential delivery phases and target deliverables")
    add_footer(slide)

    phases = _data.get("timeline_phases", [
        {"phase": "Phase 1: Inception & Discovery", "duration": "Weeks 1-4", "deliverables": "Parsed specifications, initial architectural blueprint"},
        {"phase": "Phase 2: Core Engineering & RAG Setup", "duration": "Weeks 5-12", "deliverables": "Database sync, orchestration engines integration"},
        {"phase": "Phase 3: UAT & Client Handover", "duration": "Weeks 13-16", "deliverables": "Production deployment, final document sign-off"}
    ])

    for i, phase in enumerate(phases[:3]):
        top = Inches(1.5 + (i * 1.7))
        c_shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0.5), top, Inches(3.2), Inches(1.3))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = ORANGE
        c_shape.line.fill.background()
        
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = safe_text(phase.get("phase", ""))
        p_c.alignment = PP_ALIGN.CENTER
        set_font(p_c.runs[0], size=11, bold=True, color=WHITE)
        
        p_dur = tf_c.add_paragraph()
        p_dur.text = f"({safe_text(phase.get('duration', ''))})"
        p_dur.alignment = PP_ALIGN.CENTER
        set_font(p_dur.runs[0], size=10, color=GOLD)

        d_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), top, Inches(5.5), Inches(1.3))
        d_box.fill.solid()
        d_box.fill.fore_color.rgb = OFF_WHITE
        d_box.line.color.rgb = CHARCOAL
        d_box.line.width = Pt(1)
        
        tf_d = d_box.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = "Key Deliverables / Activities:"
        set_font(p_d.runs[0], size=11, bold=True, color=CHARCOAL)
        
        p_d_desc = tf_d.add_paragraph()
        p_d_desc.text = safe_text(phase.get("deliverables", ""))
        set_font(p_d_desc.runs[0], size=10, color=CHARCOAL)
    return "Timeline slide created successfully."

@tool
def save_presentation(output_path: str) -> str:
    """Saves the presentation to the specified output path."""
    global _prs
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    _prs.save(output_path)
    return f"Presentation saved successfully at: {output_path}"

def generate_pptx_dynamically(data: dict, output_path: str):
    """
    Agentic entry point. Uses an LLM to decide which slides to add based on the input data,
    then saves the output.
    """
    global _prs, _data
    _prs = Presentation()
    _prs.slide_width = Inches(10)
    _prs.slide_height = Inches(7.5)
    _data = data

    tools = [
        create_title_slide,
        create_business_summary_slide,
        create_gap_analysis_slide,
        create_timeline_slide,
        save_presentation
    ]

    llm = ChatOpenAI(model="gpt-4o-mini", temperature=0) # Requires OPENAI_API_KEY in env
    agent_executor = create_react_agent(llm, tools)

    system_prompt = f"""
    You are an intelligent Presentation Building Agent. Your job is to construct a PPTX proposal for a client.
    You have tools to create specific slides. 
    Analyze the available data keys and decide which slides are relevant to create.
    Always start by creating the title slide. 
    Then, add other slides like business summary, gap analysis, and timeline depending on if the data seems to warrant it.
    The available data keys are: {list(data.keys())}
    
    IMPORTANT: Once you have added all necessary slides, you MUST call the `save_presentation` tool 
    with the exact output_path provided below to finalize the document.
    Output Path: {output_path}
    """

    print("Agent is dynamically deciding the presentation structure...")
    response = agent_executor.invoke({"messages": [{"role": "user", "content": system_prompt}]})
    print("Agent finished execution.")
    return output_path
